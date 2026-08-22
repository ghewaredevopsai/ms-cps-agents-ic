#!/usr/bin/env python3
"""
plan-to-pptx.py - turn a slide plan into a real PowerPoint file.

    python3 plan-to-pptx.py slide-plan-sample.json investment-report-sample.pptx

The agent produces the plan; this produces the file. In your own environment
this script is the box your existing generator already fills - the interface
between them is the plan, not the code.

    pip install python-pptx
    soffice --headless --convert-to pdf investment-report-sample.pptx   # if you need a PDF
"""
import json
import re
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x12, 0x26, 0x3F)
TEAL = RGBColor(0x1E, 0x88, 0xA8)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GREY = RGBColor(0x5C, 0x7D, 0x94)
PALE = RGBColor(0xF4, 0xF8, 0xFB)
RULE = RGBColor(0xDB, 0xE3, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xB4, 0x53, 0x09)
FONT = "Segoe UI"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
BODY_W = SLIDE_W - 2 * MARGIN
BODY_TOP = Inches(1.55)
BODY_BOT = Inches(6.7)
SLICE_COLOURS = [TEAL, NAVY, GOLD, GREY, RGBColor(0x9E, 0xC5, 0xD8), AMBER]


# ---------------------------------------------------------------- helpers


def textbox(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def line(tf, text, size, colour, bold=False, first=False, space_before=0, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.space_before = Pt(space_before)
    if align is not None:
        p.alignment = align
    f = p.font
    f.name, f.size, f.bold, f.color.rgb = FONT, Pt(size), bold, colour
    return p


def box(slide, x, y, w, h, fill, edge=None):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if edge is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = edge
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def fit(text, base, floor, per_line):
    """Shrink type rather than overflow the box.

    python-pptx cannot compute autofit - it has no text metrics - so a long
    string silently spills. Approximating here is the difference between a
    deck that survives real commentary and one that only survives the sample.
    """
    lines = max(1, -(-len(text) // per_line))
    return max(floor, base - 2 * (lines - 1))


def source(tf, text, space_before=8):
    """Every figure carries where it came from. That is the whole discipline."""
    if text:
        line(tf, text, 9, GREY, space_before=space_before).font.italic = True


# ---------------------------------------------------------------- chrome


def new_slide(prs, title, deck, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if title:
        tf = textbox(slide, MARGIN, Inches(0.42), BODY_W, Inches(0.6))
        line(tf, title, fit(title, 26, 18, 62), NAVY, bold=True, first=True)
        box(slide, MARGIN, Inches(1.2), BODY_W, Pt(2.5), TEAL)
    footer = "%s  ·  %s  ·  data as of %s" % (
        deck.get("deck_title", ""),
        deck.get("client", ""),
        deck.get("as_of", ""),
    )
    tf = textbox(slide, MARGIN, Inches(6.95), BODY_W - Inches(0.5), Inches(0.3))
    line(tf, footer, 8, GREY, first=True)
    tf = textbox(slide, SLIDE_W - MARGIN - Inches(0.5), Inches(6.95), Inches(0.5), Inches(0.3))
    line(tf, str(index), 8, GREY, first=True, align=PP_ALIGN.RIGHT)
    return slide


# ---------------------------------------------------------------- layouts


def title_slide(prs, s, deck, index):
    slide = new_slide(prs, None, deck, index)
    box(slide, Emu(0), Emu(0), SLIDE_W, Inches(4.15), NAVY)
    tf = textbox(slide, MARGIN, Inches(1.5), BODY_W, Inches(1.4))
    line(tf, s.get("title", ""), fit(s.get("title", ""), 40, 28, 46), WHITE, bold=True, first=True)
    box(slide, MARGIN, Inches(3.05), Inches(1.6), Pt(3.5), GOLD)
    tf = textbox(slide, MARGIN, Inches(3.35), BODY_W, Inches(0.6))
    line(tf, s.get("subtitle", ""), 15, RGBColor(0x9E, 0xC5, 0xD8), first=True)


def kpi_slide(prs, s, deck, index):
    slide = new_slide(prs, s.get("title", ""), deck, index)
    kpis = s.get("kpis", [])
    cols, gap = 3, Inches(0.25)
    rows = -(-len(kpis) // cols)
    cw = int((BODY_W - gap * (cols - 1)) / cols)
    ch = int((BODY_BOT - BODY_TOP - gap * (rows - 1)) / rows)
    for i, k in enumerate(kpis):
        x = MARGIN + (cw + gap) * (i % cols)
        y = BODY_TOP + (ch + gap) * (i // cols)
        box(slide, x, y, cw, ch, PALE, RULE)
        box(slide, x, y, Pt(3), ch, TEAL)
        tf = textbox(slide, x + Inches(0.22), y + Inches(0.22), cw - Inches(0.44), ch - Inches(0.44))
        line(tf, k.get("label", ""), 11, GREY, first=True)
        value = k.get("value", "")
        line(tf, value, fit(value, 26, 17, 18), NAVY, bold=True, space_before=6)
        source(tf, k.get("source", ""), space_before=6)


def numeric_column(rows, c):
    """Right-align figures, left-align words - the rule a table reads by."""
    values = [str(row[c]) for row in rows if c < len(row) and str(row[c]).strip()]
    return bool(values) and all(re.fullmatch(r"[\d.,%+-]+", v) for v in values)


def table_slide(prs, s, deck, index):
    slide = new_slide(prs, s.get("title", ""), deck, index)
    cols, rows = s.get("columns", []), s.get("rows", [])
    n = len(rows) + 1
    height = min(Inches(0.42) * n, BODY_BOT - BODY_TOP - Inches(0.6))
    shape = slide.shapes.add_table(n, len(cols), MARGIN, BODY_TOP, BODY_W, height)
    table = shape.table
    table.first_row = table.horz_banding = False
    right = [c > 0 and numeric_column(rows, c) for c in range(len(cols))]
    for c, name in enumerate(cols):
        table.cell(0, c).text = str(name)
    for r, row in enumerate(rows, start=1):
        for c, cell in enumerate(row):
            table.cell(r, c).text = str(cell)
    for r in range(n):
        for c in range(len(cols)):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (WHITE if r % 2 else PALE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.12)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if right[c] else PP_ALIGN.LEFT
            f = p.font
            f.name, f.size = FONT, Pt(12 if r == 0 else 13)
            f.bold = r == 0
            f.color.rgb = WHITE if r == 0 else NAVY
    tf = textbox(slide, MARGIN, BODY_TOP + height + Inches(0.15), BODY_W, Inches(0.4))
    source(tf, "Source: " + s["source"] if s.get("source") else "", space_before=0)


def chart_slide(prs, s, deck, index):
    slide = new_slide(prs, s.get("title", ""), deck, index)
    series = s.get("series", [])
    data = CategoryChartData()
    data.categories = [p.get("label", "") for p in series]
    data.add_series(s.get("title", "Series"), [p.get("value", 0) for p in series])
    cx, cy = Inches(8.6), Inches(4.5)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, int((SLIDE_W - cx) / 2), BODY_TOP, cx, cy, data
    )
    chart = frame.chart
    chart.font.name, chart.font.size = FONT, Pt(12)
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.number_format, labels.number_format_is_linked = '0"%"', False
    labels.position = XL_LABEL_POSITION.OUTSIDE_END
    labels.font.size, labels.font.bold, labels.font.color.rgb = Pt(12), True, NAVY
    for i, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = SLICE_COLOURS[i % len(SLICE_COLOURS)]
    tf = textbox(slide, MARGIN, BODY_TOP + cy + Inches(0.15), BODY_W, Inches(0.4))
    source(tf, "Source: " + s["source"] if s.get("source") else "", space_before=0)


def comment_slide(prs, s, deck, index):
    slide = new_slide(prs, s.get("title", ""), deck, index)
    bullets = s.get("bullets", [])
    tf = textbox(slide, MARGIN + Inches(0.1), BODY_TOP + Inches(0.15), BODY_W - Inches(0.2), Inches(4.2))
    for i, text in enumerate(bullets):
        p = line(tf, "•   " + text, fit(text, 18, 13, 88), NAVY, first=(i == 0), space_before=0 if i == 0 else 16)
        p.line_spacing = 1.15
    tf = textbox(slide, MARGIN, BODY_BOT - Inches(0.3), BODY_W, Inches(0.4))
    source(tf, "Source: " + s["source"] if s.get("source") else "", space_before=0)


def unresolved_slide(prs, items, deck, index):
    slide = new_slide(prs, "Not sourced - deliberately excluded from this deck", deck, index)
    height = Inches(0.55) * len(items) + Inches(0.9)
    box(slide, MARGIN, BODY_TOP, BODY_W, height, RGBColor(0xFF, 0xFB, 0xF1), GOLD)
    tf = textbox(slide, MARGIN + Inches(0.35), BODY_TOP + Inches(0.3), BODY_W - Inches(0.7), height)
    line(tf, "Every figure below was requested and could not be sourced.", 14, AMBER, bold=True, first=True)
    for text in items:
        line(tf, "•   " + text, fit(text, 14, 11, 110), NAVY, space_before=14)
    tf = textbox(slide, MARGIN, BODY_TOP + height + Inches(0.3), BODY_W, Inches(0.8))
    line(
        tf,
        "A gap is reported, not filled. Anything the agent could not trace to a document "
        "or an action is listed here rather than estimated onto a slide.",
        13,
        GREY,
        first=True,
    )


LAYOUTS = {
    "title": title_slide,
    "kpi": kpi_slide,
    "table": table_slide,
    "chart": chart_slide,
    "comment": comment_slide,
}


# ---------------------------------------------------------------- entry


def build(plan, out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    index = 0
    for s in plan.get("slides", []):
        layout = s.get("layout")
        if layout not in LAYOUTS:
            print("  ! skipped unknown layout: %r" % layout, file=sys.stderr)
            continue
        index += 1
        LAYOUTS[layout](prs, s, plan, index)
    if plan.get("unresolved"):
        index += 1
        unresolved_slide(prs, plan["unresolved"], plan, index)
    prs.save(out_path)
    return index


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit("usage: plan-to-pptx.py <slide-plan.json> <output.pptx>")
    with open(sys.argv[1], encoding="utf-8") as fh:
        plan = json.load(fh)
    count = build(plan, sys.argv[2])
    print("wrote %s - %d slides" % (sys.argv[2], count))
